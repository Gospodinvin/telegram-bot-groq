# pptx_export.py — экспорт в PowerPoint
import io
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def build_pptx(title: str, text: str, is_diploma: bool = False) -> io.BytesIO:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    # Титульный слайд
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Сгенерировано LectureX Bot • {datetime.now().strftime('%d.%m.%Y')}"

    if is_diploma:
        sections = re.split(r'===\s*(.+?)\s*===', text)
        for i in range(1, len(sections), 2):
            heading = sections[i].strip()
            body = sections[i+1].strip() if i+1 < len(sections) else ""
            if not heading:
                continue
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = heading
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = body[:2000]
    else:
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        current_heading = "Содержание"
        current_body = []
        for line in lines:
            if line.endswith(':') or (line.isupper() and len(line) > 3 and len(line) < 80):
                if current_body:
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = current_heading[:50]
                    content = slide.placeholders[1]
                    tf = content.text_frame
                    tf.text = '\n'.join(current_body)
                    current_body = []
                current_heading = line.rstrip(':')
            else:
                current_body.append(line)
        if current_body:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = current_heading[:50]
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = '\n'.join(current_body)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf